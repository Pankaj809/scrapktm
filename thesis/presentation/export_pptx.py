#!/usr/bin/env python3
"""Export defense slides to editable PPTX (companion to presentation.tex)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES = [
    {
        "title": "Retrieval-Augmented Generation for Nepali Government Documents",
        "subtitle": "A Hybrid Pipeline for Low-Resource, Code-Switched Legal Text",
        "bullets": [
            "ADHIKARI PANKAJ | Student ID: 2120256037",
            "Nankai University | Masters Student",
        ],
        "notes": "Open with the citizen-facing problem: laws live in PDFs, not chatbots.",
    },
    {
        "title": "Introduction: What is RAG?",
        "bullets": [
            "RAG = search official PDFs first, then let the AI read those paragraphs to answer.",
            "Kathmandu Valley bylaws are in Nepali PDFs with Nepali + English queries.",
            "General chatbots do not have your local municipal files.",
            "Goal: local, measurable retrieval for municipal legal assistance.",
        ],
        "notes": "Define RAG in plain language before any jargon.",
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Legacy PDF font encoding produces unreadable extracted text (not Unicode Nepali).",
            "Code-switched queries stress English-only embedding models.",
            "Evaluation can show 0% if ground-truth chunk IDs do not match the indexed corpus.",
            "Example: sf7df8f}F dxfgu/kflnsf instead of readable municipal text.",
        ],
        "notes": "Stress encoding + evaluation hygiene, not only model failure.",
    },
    {
        "title": "Dataset & Pipeline Overview",
        "bullets": [
            "400+ PDFs scraped -> 50 curated gold documents -> 446 JSONL chunks.",
            "Sources: Kathmandu, Lalitpur, Bhaktapur + federal laws.",
            "Pipeline: scrape -> curate -> PyMuPDF/Surya OCR -> chunk -> retrieve -> answer.",
            "830 pages total; 45 text-based + 5 scanned PDFs.",
        ],
        "notes": "Point to curated_corpus_manifest.jsonl and extraction_output/chunks.",
    },
    {
        "title": "RAG Architecture",
        "bullets": [
            "User query -> vectorize (TF-IDF or dense) -> search index -> top-k chunks -> LLM answer.",
            "Offline index built from: PDFs -> extract -> chunk -> local_rag_db/documents.jsonl.",
            "local_rag_db = project-owned local retrieval database (on-disk chunk store).",
            "Retrieval quality limits answer quality: miss the chunk, risk hallucination.",
        ],
        "notes": "Walk left-to-right on the diagram if presenting with PDF.",
    },
    {
        "title": "Methodology",
        "bullets": [
            "Router: text PDFs -> PyMuPDF; scanned -> Surya OCR (Apple MPS).",
            "Chunking: heading-aware sections; 512-token fallback with overlap.",
            "Dense baseline: BGE-M3 + all-MiniLM-L6-v2 with FAISS.",
            "Hybrid fix: metadata aliasing + scikit-learn TF-IDF (sparse keyword retrieval).",
            "BM25 is same family as TF-IDF; our code implements TF-IDF.",
        ],
        "notes": "Tools: Python, LangChain Core, scikit-learn, FAISS, sentence-transformers.",
    },
    {
        "title": "Literature Review (Brief)",
        "bullets": [
            "Low-resource NLP: English models underperform on Devanagari + code-switching.",
            "South Asian PDFs: legacy fonts break standard OCR and embeddings.",
            "Hybrid retrieval: sparse + dense improves robustness on noisy text.",
            "Legal IR needs exact fees, tables, and document lists.",
        ],
        "notes": "Keep short for a 5-minute defense.",
    },
    {
        "title": "Aim & Objectives",
        "bullets": [
            "Aim: local RAG that retrieves correct clauses from Nepali municipal PDFs.",
            "Curate reproducible corpus (50 PDFs -> JSONL chunks).",
            "Measure dense baselines on raw extracted text.",
            "Design alias-augmented sparse retrieval; compare Recall@k and MRR.",
            "Contrast grounded RAG vs direct LLM guesses on tax/fee queries.",
        ],
        "notes": None,
    },
    {
        "title": "Understanding the Metrics",
        "bullets": [
            "Recall@k: was the correct chunk in the top k results? (1=yes, 0=no).",
            "MRR: 1/rank of first correct chunk (rank 1 -> 1.0, rank 2 -> 0.5).",
            "Critical for RAG: if retrieval fails, the LLM hallucinates.",
            "We report Recall@3 for hybrid system; Recall@5 for dense FAISS benchmark.",
        ],
        "notes": "Use a concrete example: top-3 paragraphs shown to the LLM.",
    },
    {
        "title": "Results: Why Raw Retrieval Fails",
        "bullets": [
            "Sparse TF-IDF on raw corrupted text: Recall@3 = 0.091, MRR = 0.030 (1/11 hits).",
            "Cause: 26/50 documents extract as Preeti font gibberish.",
            "Dense models (BGE-M3, MiniLM) exceed the offline memory budget -> deferred, not reported as zeros.",
            "The failure is font-encoding, not the retrieval model.",
        ],
        "notes": "Numbers from paper_analysis.py over the real corpus.",
    },
    {
        "title": "Results: Retrieval Performance",
        "bullets": [
            "Raw corrupted text: Recall@3 = 0.091 | MRR = 0.030 | DocR@3 = 0.091",
            "Deterministic Preeti repair (no aliases): Recall@3 = 0.364 | MRR = 0.197 | DocR@3 = 0.455",
            "Metadata aliasing: Recall@3 = 1.000 -- ORACLE (query leakage), not an operational score.",
            "Key result: deterministic repair gives a 4-5x gain with no hand-authored knowledge.",
        ],
        "notes": "Repair is the contribution; aliasing perfect score is leakage.",
    },
    {
        "title": "Key Finding: Direct LLM vs Local RAG",
        "bullets": [
            "Parking fee query: ChatGPT guesses ~Rs 20/hr; RAG returns KTM_FIN_002 clause.",
            "Tax slab query: LLM lacks exact local schedule; RAG returns precise rate.",
            "Map-pass documents: LLM gives generic list; RAG returns full bylaw checklist.",
            "RAG gives citable, numerically exact values from source PDFs.",
        ],
        "notes": "Emphasize grounding and citations for citizen trust.",
    },
    {
        "title": "Conclusion & Future Work",
        "bullets": [
            "Built scrape -> curate -> extract -> repair -> chunk -> retrieve pipeline.",
            "Raw retrieval fails on corrupted Nepali PDFs (Recall@3 = 0.091).",
            "Deterministic Preeti -> Unicode repair lifts Recall@3 to 0.364 (doc-level 0.455), no aliases.",
            "Next: learned transliteration repair, BM25+dense fusion, larger human-labeled queries.",
        ],
        "notes": "Close on accessible legal information for Kathmandu Valley.",
    },
    {
        "title": "Thank You",
        "bullets": ["Questions?", "Repository: scrapktm/thesis/"],
        "notes": "Backup: aliases are a research bridge until proper Unicode OCR is deployed.",
    },
]


def add_title_slide(prs: Presentation, slide_data: dict) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = slide_data["title"]
    if slide.placeholders[1].text_frame:
        sub = slide_data.get("subtitle", "")
        body = "\n".join(slide_data.get("bullets", []))
        slide.placeholders[1].text = f"{sub}\n{body}".strip()
    if slide_data.get("notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def add_content_slide(prs: Presentation, slide_data: dict) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = slide_data["title"]
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(slide_data.get("bullets", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
    if slide_data.get("notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def main() -> None:
    out = Path(__file__).with_name("presentation.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(SLIDES):
        if i == 0:
            add_title_slide(prs, data)
        else:
            add_content_slide(prs, data)

    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
